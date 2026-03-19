"""Pure synchronous functions for .pptx format manipulation.

No network calls, no tokens. Takes bytes in, returns bytes out.
Independently testable by creating real .pptx content in memory.
"""

from __future__ import annotations

import io

from pptx import Presentation as NewPresentation
from pptx.presentation import Presentation
from pptx.util import Inches

# Default layout index mapping for the standard python-pptx template.
_LAYOUT_PRESETS: dict[str, int] = {
    "blank": 6,
    "title": 0,
    "title_and_content": 1,
    "section_header": 2,
    "two_content": 3,
    "title_only": 5,
}


def resolve_layout_index(prs: Presentation, layout: str) -> int:
    """Resolve a layout name to an index in the presentation's slide layouts.

    Tries an exact name match (case-insensitive, underscores/hyphens/spaces
    normalized) against the layouts that actually exist in the file first.
    Falls back to the preset map, then to index 0.

    Args:
        prs: An open Presentation object.
        layout: Layout name string (e.g. "title_and_content", "blank").

    Returns:
        The integer index into prs.slide_layouts.
    """
    normalized = layout.lower().replace(" ", "_").replace("-", "_")

    for idx, sl in enumerate(prs.slide_layouts):
        sl_normalized = sl.name.lower().replace(" ", "_").replace("-", "_")
        if sl_normalized == normalized:
            return idx

    if normalized in _LAYOUT_PRESETS:
        preset_idx = _LAYOUT_PRESETS[normalized]
        if preset_idx < len(prs.slide_layouts):
            return preset_idx

    return 0


def build_pptx(title: str = "") -> bytes:
    """Create a new presentation with an optional title slide.

    Args:
        title: If non-empty, creates a Title Slide layout with this text.
            Otherwise, creates a single blank slide.

    Returns:
        The .pptx file content as bytes.
    """
    prs = NewPresentation()

    if title:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = title
    else:
        slide_layout = prs.slide_layouts[6]
        prs.slides.add_slide(slide_layout)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def extract_slides(
    pptx_bytes: bytes,
    include_notes: bool = False,
) -> list[dict]:
    """Extract structured text data from every slide in a presentation.

    Args:
        pptx_bytes: Raw .pptx file bytes.
        include_notes: Whether to include speaker notes.

    Returns:
        A list of dicts, one per slide, each containing:
        number, title, texts (list of str), and notes (str).
    """
    prs = NewPresentation(io.BytesIO(pptx_bytes))
    slides: list[dict] = []

    for i, slide in enumerate(prs.slides):
        slide_data: dict = {
            "number": i + 1,
            "title": "",
            "texts": [],
            "notes": "",
        }

        title_shape = slide.shapes.title
        if title_shape and title_shape.has_text_frame:
            slide_data["title"] = title_shape.text_frame.text.strip()

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and text != slide_data["title"]:
                    slide_data["texts"].append(text)

        if include_notes and slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_data["notes"] = notes_text

        slides.append(slide_data)

    return slides


def add_slide(
    pptx_bytes: bytes,
    layout: str,
    title: str,
    content: str,
) -> tuple[bytes, int]:
    """Add a slide to an existing presentation.

    Args:
        pptx_bytes: Raw .pptx file bytes.
        layout: Layout name (e.g. "blank", "title_and_content").
        title: Title text for the new slide (if layout supports it).
        content: Body text for the new slide (if layout supports it).

    Returns:
        A tuple of (updated .pptx bytes, total slide count after insertion).
    """
    prs = NewPresentation(io.BytesIO(pptx_bytes))

    layout_idx = resolve_layout_index(prs, layout)
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)

    # Placeholder idx 0 = title, idx 1 = body in the standard template.
    # A more robust approach would use ph.placeholder_format.type
    # (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.BODY) which is template-independent,
    # but not all custom templates set placeholder types correctly.
    # Using numeric indices for now as it works reliably with standard layouts.
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0 and title:
            ph.text = title
        elif idx == 1 and content:
            ph.text = content

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), len(prs.slides)


def update_text(
    pptx_bytes: bytes,
    slide_number: int,
    text: str,
    shape_index: int,
) -> tuple[bytes, str, str | None]:
    """Update text in a shape on the given slide.

    If the target slide has no text shapes, a new text box is created.

    Args:
        pptx_bytes: Raw .pptx file bytes.
        slide_number: 1-based slide number.
        text: New text content.
        shape_index: 0-based index among text-bearing shapes on the slide.

    Returns:
        A tuple of (updated bytes, shape name, error string or None).
    """
    prs = NewPresentation(io.BytesIO(pptx_bytes))

    if slide_number < 1 or slide_number > len(prs.slides):
        return (
            pptx_bytes,
            "",
            f"Slide number {slide_number} is out of range (1-{len(prs.slides)})",
        )

    slide = prs.slides[slide_number - 1]
    text_shapes = [s for s in slide.shapes if s.has_text_frame]

    if not text_shapes:
        txbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        txbox.text_frame.text = text
        shape_name = txbox.name
    elif shape_index < 0 or shape_index >= len(text_shapes):
        return (
            pptx_bytes,
            "",
            (
                f"Shape index {shape_index} is out of range "
                f"(0-{len(text_shapes) - 1}). "
                f"Slide {slide_number} has {len(text_shapes)} text shape(s)."
            ),
        )
    else:
        shape = text_shapes[shape_index]
        shape.text_frame.text = text
        shape_name = shape.name

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), shape_name, None
