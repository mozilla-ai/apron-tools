"""Pure unit tests for .pptx format manipulation.

No mocks — creates real .pptx bytes and verifies their content.
"""

from __future__ import annotations

import io

from pptx import Presentation

from apron_tools.providers.microsoft.powerpoint.presentation import (
    add_slide,
    build_pptx,
    extract_slides,
    resolve_layout_index,
    update_text,
)

# ---------------------------------------------------------------------------
# build_pptx
# ---------------------------------------------------------------------------


class TestBuildPptx:
    def test_blank(self) -> None:
        data = build_pptx()
        prs = Presentation(io.BytesIO(data))
        assert len(prs.slides) == 1

    def test_with_title(self) -> None:
        data = build_pptx("My Presentation")
        prs = Presentation(io.BytesIO(data))
        assert len(prs.slides) == 1
        title_shape = prs.slides[0].shapes.title
        assert title_shape is not None
        assert title_shape.text == "My Presentation"

    def test_empty_title_creates_blank(self) -> None:
        data = build_pptx("")
        prs = Presentation(io.BytesIO(data))
        assert len(prs.slides) == 1


# ---------------------------------------------------------------------------
# extract_slides
# ---------------------------------------------------------------------------


class TestExtractSlides:
    def test_basic_extraction(self) -> None:
        data = build_pptx("Hello World")
        slides = extract_slides(data)
        assert len(slides) == 1
        assert slides[0]["number"] == 1
        assert slides[0]["title"] == "Hello World"

    def test_blank_slide_extraction(self) -> None:
        data = build_pptx()
        slides = extract_slides(data)
        assert len(slides) == 1
        assert slides[0]["title"] == ""
        assert slides[0]["texts"] == []

    def test_notes_excluded_by_default(self) -> None:
        data = build_pptx("Title")
        slides = extract_slides(data, include_notes=False)
        assert slides[0]["notes"] == ""

    def test_notes_included_when_requested(self) -> None:
        """Build a presentation with speaker notes and verify extraction."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "With Notes"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Remember this point"
        buf = io.BytesIO()
        prs.save(buf)

        slides = extract_slides(buf.getvalue(), include_notes=True)
        assert slides[0]["notes"] == "Remember this point"

    def test_multi_slide_extraction(self) -> None:
        data = build_pptx("Slide 1")
        data, _ = add_slide(data, "title_and_content", "Slide 2", "Body text")
        slides = extract_slides(data)
        assert len(slides) == 2
        assert slides[0]["title"] == "Slide 1"
        assert slides[1]["title"] == "Slide 2"


# ---------------------------------------------------------------------------
# add_slide
# ---------------------------------------------------------------------------


class TestAddSlide:
    def test_add_blank_slide(self) -> None:
        data = build_pptx("First")
        updated, count = add_slide(data, "blank", "", "")
        assert count == 2
        prs = Presentation(io.BytesIO(updated))
        assert len(prs.slides) == 2

    def test_add_title_and_content(self) -> None:
        data = build_pptx()
        updated, count = add_slide(data, "title_and_content", "New Title", "New Content")
        assert count == 2
        slides = extract_slides(updated)
        assert slides[1]["title"] == "New Title"
        assert "New Content" in slides[1]["texts"]

    def test_add_multiple_slides(self) -> None:
        data = build_pptx("First")
        data, count = add_slide(data, "blank", "", "")
        assert count == 2
        data, count = add_slide(data, "title", "Third", "")
        assert count == 3
        prs = Presentation(io.BytesIO(data))
        assert len(prs.slides) == 3

    def test_returns_updated_bytes(self) -> None:
        original = build_pptx()
        updated, _ = add_slide(original, "blank", "", "")
        assert original != updated


# ---------------------------------------------------------------------------
# update_text
# ---------------------------------------------------------------------------


class TestUpdateText:
    def test_update_existing_shape(self) -> None:
        data = build_pptx("Original Title")
        updated, shape_name, error = update_text(data, 1, "New Title", 0)
        assert error is None
        assert shape_name != ""
        slides = extract_slides(updated)
        assert slides[0]["title"] == "New Title"

    def test_slide_out_of_range(self) -> None:
        data = build_pptx()
        _, _, error = update_text(data, 5, "text", 0)
        assert error is not None
        assert "out of range" in error

    def test_slide_zero_out_of_range(self) -> None:
        data = build_pptx()
        _, _, error = update_text(data, 0, "text", 0)
        assert error is not None
        assert "out of range" in error

    def test_shape_index_out_of_range(self) -> None:
        data = build_pptx("Title")
        _, _, error = update_text(data, 1, "text", 99)
        assert error is not None
        assert "Shape index" in error

    def test_creates_textbox_when_no_text_shapes(self) -> None:
        """On a blank slide with no text shapes, a new textbox is created."""
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]
        prs.slides.add_slide(slide_layout)
        # Remove all shapes to guarantee none have text frames.
        slide = prs.slides[0]
        for shape in list(slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)
        buf = io.BytesIO()
        prs.save(buf)

        updated, shape_name, error = update_text(buf.getvalue(), 1, "New text", 0)
        assert error is None
        assert shape_name != ""
        slides = extract_slides(updated)
        assert "New text" in slides[0]["texts"]


# ---------------------------------------------------------------------------
# resolve_layout_index
# ---------------------------------------------------------------------------


class TestResolveLayoutIndex:
    def test_preset_blank(self) -> None:
        prs = Presentation()
        idx = resolve_layout_index(prs, "blank")
        assert idx == 6

    def test_preset_title(self) -> None:
        prs = Presentation()
        idx = resolve_layout_index(prs, "title")
        assert idx == 0

    def test_preset_title_and_content(self) -> None:
        prs = Presentation()
        idx = resolve_layout_index(prs, "title_and_content")
        assert idx == 1

    def test_case_insensitive(self) -> None:
        prs = Presentation()
        idx = resolve_layout_index(prs, "BLANK")
        assert idx == 6

    def test_hyphen_normalisation(self) -> None:
        prs = Presentation()
        idx = resolve_layout_index(prs, "title-and-content")
        assert idx == 1

    def test_space_normalisation(self) -> None:
        prs = Presentation()
        idx = resolve_layout_index(prs, "title and content")
        assert idx == 1

    def test_unknown_layout_falls_back_to_zero(self) -> None:
        prs = Presentation()
        idx = resolve_layout_index(prs, "nonexistent_layout")
        assert idx == 0
